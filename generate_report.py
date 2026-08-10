#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
시장조사 보고서 템플릿 v2.0 자동 채우기 스크립트
(260810_시장조사_양식_개선_v_2_0.pptx 기준)

사용법:
    python3 generate_report_v2.py data.json output.pptx
"""
import os
import sys
import json
import re
import zipfile
import shutil
from collections import defaultdict
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "template_v2.pptx")

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

JUDGE_COLORS = {"양호": "2E9E4F", "보통": "E8A93B", "주의": "C0392B"}
GRADE_BLUE = "0070C0"


def find_by_id(shapes, sid):
    for shp in shapes:
        if shp.shape_id == sid:
            return shp
        if shp.shape_type == 6:  # GROUP
            r = find_by_id(shp.shapes, sid)
            if r is not None:
                return r
    return None


def clean_text(s):
    if not isinstance(s, str):
        return s
    return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", s)


def set_para_text(paragraph, new_text):
    runs = paragraph.runs
    if not runs:
        return
    runs[0].text = clean_text(new_text)
    for r in runs[1:]:
        r.text = ""


def set_variable_paragraphs(shape, texts):
    """텍스트박스 안의 문단 개수를 입력한 texts 개수에 맞게 늘리거나 줄인다.
    (예: 종합의견이 2줄이든 5줄이든 고정된 3줄에 얽매이지 않고 유연하게 대응)"""
    import copy
    tf = shape.text_frame
    txBody = tf._txBody
    paras_el = txBody.findall(f"{{{A_NS}}}p")
    n_have, n_need = len(paras_el), len(texts)
    if n_need == 0:
        n_need = 1
        texts = [""]
    if n_need > n_have:
        last = paras_el[-1]
        for _ in range(n_need - n_have):
            txBody.append(copy.deepcopy(last))
    elif n_need < n_have:
        for p in paras_el[n_need:]:
            txBody.remove(p)
    for p, text in zip(tf.paragraphs, texts):
        set_para_text(p, text)


def set_shape_text(shape, new_text, para_idx=0):
    set_para_text(shape.text_frame.paragraphs[para_idx], new_text)


def set_run_color(run, hex_color):
    run.font.color.rgb = RGBColor.from_string(hex_color)


def set_oval_style(shape, is_blue):
    """등급 배지(oval)의 채우기/선/글자색을 S,A(파란색) 또는 B,C,D(연한색) 스타일로 전환."""
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sp = shape._element
    spPr = sp.find(f"{{{P_NS}}}spPr")
    solidFill = spPr.find(f"{{{A_NS}}}solidFill")
    if solidFill is not None:
        spPr.remove(solidFill)
    new_fill = etree.SubElement(spPr, f"{{{A_NS}}}solidFill")
    # solidFill은 prstGeom 다음에 와야 하므로 위치 재조정
    prstGeom = spPr.find(f"{{{A_NS}}}prstGeom")
    spPr.remove(new_fill)
    prstGeom.addnext(new_fill)

    if is_blue:
        clr = etree.SubElement(new_fill, f"{{{A_NS}}}srgbClr")
        clr.set("val", GRADE_BLUE)
        text_scheme = "bg1"
    else:
        clr = etree.SubElement(new_fill, f"{{{A_NS}}}schemeClr")
        clr.set("val", "accent5")
        lm = etree.SubElement(clr, f"{{{A_NS}}}lumMod")
        lm.set("val", "20000")
        lo = etree.SubElement(clr, f"{{{A_NS}}}lumOff")
        lo.set("val", "80000")
        text_scheme = "tx1"

    for run in shape.text_frame.paragraphs[0].runs:
        rPr = run._r.find(f"{{{A_NS}}}rPr")
        if rPr is not None:
            sf = rPr.find(f"{{{A_NS}}}solidFill")
            if sf is not None:
                rPr.remove(sf)
            new_sf = etree.SubElement(rPr, f"{{{A_NS}}}solidFill")
            sc = etree.SubElement(new_sf, f"{{{A_NS}}}schemeClr")
            sc.set("val", text_scheme)
            # solidFill은 rPr의 첫 자식이어야 함
            rPr.remove(new_sf)
            rPr.insert(0, new_sf)


def fill_text(slide, d):
    S = lambda sid: find_by_id(slide.shapes, sid)

    # 제목 / 작성일
    title = d["title"].strip()
    suffix = "시장조사 보고서"
    if title.endswith(suffix):
        title = title[: -len(suffix)].strip()
    set_shape_text(S(3), f"{title} {suffix}")
    set_shape_text(S(4), f"’{d['date']}  ㅣ  건축·주택마케팅팀")

    # 지역등급 뱃지
    set_shape_text(S(70), d["region_grade"])

    # 적정분양가
    set_shape_text(S(90), f"@{d['compare1_price_py']}")
    set_shape_text(S(92), f"@{d['compare2_price_py']}")
    set_shape_text(S(94), f"@{d['site_price_py']}")
    set_shape_text(S(96), f"{d['site_price_eok']}억  (@{d['site_price_py']}만원)")
    set_shape_text(S(98), f"84㎡ ({d['pyeong']}평)기준 (확장포함 · 이자후불제)")

    # 입지 코멘트 (교통/생활/개발)
    box = S(67)
    set_para_text(box.text_frame.paragraphs[0], f"교통\t{d['loc_transport']}")
    set_para_text(box.text_frame.paragraphs[1], f"생활\t{d['loc_life']}")
    set_para_text(box.text_frame.paragraphs[2], f"개발\t{d['loc_dev']}")

    # 상단 개요 요약 (2문단)
    box = S(46)
    set_para_text(box.text_frame.paragraphs[0], d["region_summary_1"])
    set_para_text(box.text_frame.paragraphs[1], d["region_summary_2"])

    # 지역위계 본문 (1문단)
    set_shape_text(S(64), d["hierarchy_note"])

    # 비교단지 라벨 (최대 5개 슬롯)
    label_ids = [79, 80, 81, 82, 83]
    labels = d["stock_labels"]
    for i, sid in enumerate(label_ids):
        if i < len(labels):
            item = labels[i]
            set_shape_text(S(sid), f"{item['name']} [{item['built']}, {item['units']}세대]")
        else:
            for p in S(sid).text_frame.paragraphs:
                set_para_text(p, "")

    # 등급 배지 5개 (수급/가격/거래/분양/미분양) - oval id, 색상은 S/A=파랑, 나머지=연한색
    grades = d["grade_table"]
    grade_oval_ids = {"supply": 87, "price": 88, "deal": 89, "presale": 91, "unsold": 93}
    for key, sid in grade_oval_ids.items():
        val = grades[key]
        shp = S(sid)
        set_shape_text(shp, val)
        set_oval_style(shp, is_blue=val in ("S", "A"))

    # 검토의견 표 (입지및시장환경/상품/브랜드영향/분양등급)
    tbl = S(48).table
    judge = d["judge_table"]  # dict: location, product, brand
    order = ["location", "product", "brand"]
    for col, key in enumerate(order):
        val = judge[key]
        run = tbl.cell(1, col).text_frame.paragraphs[0].runs[0]
        run.text = clean_text(f"● {val}")
        color_hex = JUDGE_COLORS.get(val)
        if color_hex:
            set_run_color(run, color_hex)
    grade_para = tbl.cell(1, 3).text_frame.paragraphs[0]
    grade_para.runs[0].text = clean_text(f"{d['final_grade']} 등급({d['final_grade_note']})")
    for extra_run in grade_para.runs[1:]:
        extra_run.text = ""

    # 검토의견 본문 (개수 가변 - 2줄이든 5줄이든 입력한 만큼 자동 대응)
    opinions = d["opinions"] if d.get("opinions") else [
        v for v in [d.get("opinion_1"), d.get("opinion_2"), d.get("opinion_3")] if v
    ]
    set_variable_paragraphs(S(78), opinions)


def fill_charts(slide, d):
    S = lambda sid: find_by_id(slide.shapes, sid)

    def compare_category(name, built):
        return f"{name}\n('{built}입주)" if built else name

    chart = S(11).chart
    cd = CategoryChartData()
    cd.categories = [
        compare_category(d["compare1_name"], d.get("compare1_built", "")),
        compare_category(d["compare2_name"], d.get("compare2_built", "")),
        "SITE 적정가\n(신규)",
    ]
    cd.add_series("가격", (d["compare1_price_eok"], d["compare2_price_eok"], d["site_price_eok"]))
    chart.replace_data(cd)

    chart = S(77).chart
    cd = CategoryChartData()
    cd.categories = ["입지(35)", "브랜드(15)", "단지(40)", "가격수준(10)", "연식(50)"]
    r = d["radar"]
    cd.add_series(d["compare1_name"], tuple(r["compare1"]))
    cd.add_series(d["compare2_name"], tuple(r["compare2"]))
    cd.add_series("SITE", tuple(r["site"]))
    chart.replace_data(cd)


def get_chart_part_for_shape(pptx_path, shape_id):
    """슬라이드에서 특정 shape_id의 그래픽프레임이 참조하는 차트 xml 파일명을 찾는다."""
    with zipfile.ZipFile(pptx_path) as z:
        slide_xml = z.read("ppt/slides/slide1.xml")
        rels_xml = z.read("ppt/slides/_rels/slide1.xml.rels")

    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main",
          "c": C_NS, "r": R_NS}
    tree = etree.fromstring(slide_xml)
    target_rid = None
    for gf in tree.iter(f"{{{ns['p']}}}graphicFrame"):
        cNvPr = gf.find(".//p:nvGraphicFramePr/p:cNvPr", ns)
        if cNvPr is not None and cNvPr.get("id") == str(shape_id):
            chart_el = gf.find(f".//{{{C_NS}}}chart")
            if chart_el is not None:
                target_rid = chart_el.get(f"{{{R_NS}}}id")
            break
    rels_tree = etree.fromstring(rels_xml)
    for rel in rels_tree:
        if rel.get("Id") == target_rid:
            target = rel.get("Target")
            return "ppt/" + target.replace("../", "")
    raise ValueError(f"shape {shape_id}의 차트를 찾을 수 없습니다.")


def fill_stock_chart(pptx_path, d):
    chart_rel_path = get_chart_part_for_shape(pptx_path, 60)
    chart_dir = os.path.dirname(chart_rel_path)  # 예: ppt/charts
    chart_basename = os.path.basename(chart_rel_path)  # 예: chart1.xml

    tmp_dir = "unpacked_stock_tmp_v2"
    shutil.rmtree(tmp_dir, ignore_errors=True)
    with zipfile.ZipFile(pptx_path) as z:
        z.extractall(tmp_dir)

    chart_path = os.path.join(tmp_dir, chart_rel_path)

    dongs = d["stock_dongs"]
    series = d["stock_series"]
    n = len(dongs)
    if n < 3 or n > 5:
        raise ValueError("고가/평균/저가 비교는 3~5개 동만 지원합니다 (원본 템플릿 레이아웃 제약).")

    NS = {"c": C_NS}
    tree = etree.parse(chart_path)
    root = tree.getroot()
    sers = root.findall(".//c:ser", NS)
    order = ["high", "low", "avg"]
    last_row = n + 1  # 데이터는 2행부터 시작하므로 마지막 행 = n+1

    def rebuild_pts(cache_el, values):
        pt_count = cache_el.find("c:ptCount", NS)
        pt_count.set("val", str(len(values)))
        for pt in cache_el.findall("c:pt", NS):
            cache_el.remove(pt)
        for i, val in enumerate(values):
            pt = etree.SubElement(cache_el, f"{{{C_NS}}}pt")
            pt.set("idx", str(i))
            v_el = etree.SubElement(pt, f"{{{C_NS}}}v")
            v_el.text = str(val)

    def fix_formula_range(ref_el):
        """Sheet1!$A$2:$A$6 같은 셀 범위 참조를 실제 데이터 개수(n)에 맞게 끝 행 번호 수정."""
        if ref_el is None or not ref_el.text:
            return
        ref_el.text = re.sub(r"(\$[A-Z]+\$2:\$[A-Z]+\$)\d+", rf"\g<1>{last_row}", ref_el.text)

    for ser, key in zip(sers, order):
        cat_cache_parent = ser.find(".//c:cat/c:strRef", NS)
        fix_formula_range(cat_cache_parent.find("c:f", NS) if cat_cache_parent is not None else None)
        cat_cache = ser.find(".//c:cat//c:strCache", NS)
        rebuild_pts(cat_cache, dongs)

        val_ref = ser.find(".//c:val/c:numRef", NS)
        fix_formula_range(val_ref.find("c:f", NS) if val_ref is not None else None)
        num_cache = ser.find(".//c:val//c:numCache", NS)
        rebuild_pts(num_cache, series[key])

    all_vals = list(series["high"]) + list(series["low"]) + list(series["avg"])
    data_min, data_max = min(all_vals), max(all_vals)
    pad = max((data_max - data_min) * 0.2, 1)
    new_min = max(0, int(data_min - pad))
    new_max = int(data_max + pad) + 1
    unit = max(1, round((new_max - new_min) / 5))

    val_ax = root.find(".//c:valAx", NS)
    scaling = val_ax.find("c:scaling", NS)
    for tag in ("c:max", "c:min"):
        el = scaling.find(tag, NS)
        if el is not None:
            scaling.remove(el)
    max_el = etree.SubElement(scaling, f"{{{C_NS}}}max")
    max_el.set("val", str(new_max))
    min_el = etree.SubElement(scaling, f"{{{C_NS}}}min")
    min_el.set("val", str(new_min))
    orientation = scaling.find("c:orientation", NS)
    scaling.remove(max_el)
    scaling.remove(min_el)
    orientation.addnext(max_el)
    max_el.addnext(min_el)

    major_unit = val_ax.find("c:majorUnit", NS)
    if major_unit is not None:
        major_unit.set("val", str(unit))

    tree.write(chart_path, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 임베디드 엑셀 데이터도 함께 갱신 (실제 PowerPoint가 "데이터 편집" 시 참조하는 원본)
    # - 캐시(numCache/strCache)만 바꾸고 이 워크북을 안 바꾸면, 실제 PowerPoint가 차트를
    #   다시 검증할 때 범위/데이터 불일치로 깨질 수 있어 반드시 함께 갱신해야 함
    try:
        import openpyxl
        rels_path = os.path.join(chart_dir, "_rels", chart_basename + ".rels")
        rels_tree = etree.parse(os.path.join(tmp_dir, rels_path))
        xlsx_target = None
        for rel in rels_tree.getroot():
            if rel.get("Type", "").endswith("/package"):
                xlsx_target = rel.get("Target")
                break
        if xlsx_target:
            xlsx_path = os.path.normpath(os.path.join(tmp_dir, chart_dir, xlsx_target))
            wb = openpyxl.load_workbook(xlsx_path)
            ws = wb["Sheet1"] if "Sheet1" in wb.sheetnames else wb.active
            # 기존 데이터 영역 비우기 (최대 10행까지 여유있게 클리어)
            for row in range(2, 12):
                for col in ("A", "B", "C", "D"):
                    ws[f"{col}{row}"] = None
            ws["B1"], ws["C1"], ws["D1"] = "고가", "저가", "평균가"
            for i, name in enumerate(dongs):
                r = i + 2
                ws[f"A{r}"] = name
                ws[f"B{r}"] = series["high"][i]
                ws[f"C{r}"] = series["low"][i]
                ws[f"D{r}"] = series["avg"][i]
            wb.save(xlsx_path)
    except Exception as e:
        print(f"경고: 임베디드 엑셀 데이터 갱신 실패 (차트 자체는 정상 작동): {e}")

    out_tmp = pptx_path + ".tmp"
    if os.path.exists(out_tmp):
        os.remove(out_tmp)
    with zipfile.ZipFile(out_tmp, "w", zipfile.ZIP_DEFLATED) as zf:
        for root_dir, _, files in os.walk(tmp_dir):
            for file in files:
                full = os.path.join(root_dir, file)
                rel = os.path.relpath(full, tmp_dir)
                zf.write(full, rel)
    shutil.move(out_tmp, pptx_path)
    shutil.rmtree(tmp_dir, ignore_errors=True)


def main():
    if len(sys.argv) != 3:
        print("사용법: python3 generate_report_v2.py data.json output.pptx")
        sys.exit(1)
    data_path, out_path = sys.argv[1], sys.argv[2]
    d = json.load(open(data_path, encoding="utf-8"))
    d = defaultdict(str, d)

    prs = Presentation(TEMPLATE)
    slide = prs.slides[0]
    fill_text(slide, d)
    fill_charts(slide, d)
    prs.save(out_path)

    fill_stock_chart(out_path, d)
    print(f"완료: {out_path}")


if __name__ == "__main__":
    main()
